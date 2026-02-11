#!/usr/bin/env python3
"""Convert PDFs and PPTX files in a subject folder to .txt files in documents/txt/."""

import sys
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import fitz  # PyMuPDF
from pptx import Presentation


def extract_text_from_pdf(pdf_path: Path) -> str:
    """Extract text from a PDF using PyMuPDF, with page number markers."""
    doc = fitz.open(pdf_path)
    parts = []
    for i, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            parts.append(f"--- Page {i} ---\n{text}")
        else:
            parts.append(f"--- Page {i} ---")
    doc.close()
    return "\n\n".join(parts)


def extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract text from a PowerPoint file."""
    prs = Presentation(pptx_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text)


def main():
    subjects_root = project_root / "data" / "subjects"
    if not subjects_root.exists():
        print("No data/subjects directory found.")
        return

    # Process each subject directory
    subject_dirs = sorted(d for d in subjects_root.iterdir() if d.is_dir())
    if not subject_dirs:
        print("No subject directories found in", subjects_root)
        return

    for subject_dir in subject_dirs:
        txt_dir = subject_dir / "documents" / "txt"
        txt_dir.mkdir(parents=True, exist_ok=True)

        # Collect all PDFs and PPTX under this subject (including subfolders)
        files = sorted(
            f for f in subject_dir.rglob("*")
            if f.is_file() and f.suffix.lower() in (".pdf", ".pptx")
        )

        if not files:
            print(f"⏭️  No PDF/PPTX in {subject_dir.name}")
            continue

        print(f"\n📁 {subject_dir.name}")
        for file_path in files:
            stem = file_path.stem
            txt_path = txt_dir / f"{stem}.txt"

            ext = file_path.suffix.lower()
            print(f"  📄 Converting: {file_path.name} ...", end=" ", flush=True)

            try:
                if ext == ".pdf":
                    text = extract_text_from_pdf(file_path)
                elif ext == ".pptx":
                    text = extract_text_from_pptx(file_path)
                else:
                    continue

                if not text.strip():
                    print("⚠️  No text extracted!")
                    continue

                txt_path.write_text(text, encoding="utf-8")
                word_count = len(text.split())
                print(f"✅ {word_count:,} words → {txt_path.name}")
            except Exception as e:
                print(f"❌ Error: {e}")

    print("\nDone! Txt files written under each subject's documents/txt/")


if __name__ == "__main__":
    main()
